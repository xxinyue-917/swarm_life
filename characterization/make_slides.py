#!/usr/bin/env python3
"""
Generate PowerPoint slides from characterization results.
- Grid overview slides with 5x5 screenshot grids
- 5x5 video grid slides (embedded videos)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


RESULTS_DIR = Path(__file__).parent / 'results'
GRIDS_DIR = Path(__file__).parent / 'results' / 'grids'
OUTPUT_DIR = Path(__file__).parent
GRID_VALUES = [-1.0, -0.5, 0.0, 0.5, 1.0]

SLIDE_W = 13.33  # Widescreen width in inches
SLIDE_H = 7.5


# ============================================================
# Sweep definitions
# ============================================================

SWEEPS = [
    {
        'folder': 'kpos_offdiag',
        'title': 'Sweep 1: Cross-Species Position Coupling',
        'subtitle': 'K₁₂ vs K₂₁ (K₁₁=K₂₂=0.6, K_rot=0)',
        'p1': 'K₁₂', 'p2': 'K₂₁',
        'grid': 'kpos_offdiag.png',
        'cases': None,
    },
    {
        'folder': 'kpos_x_krot',
        'title': 'Sweep 2: Position × Rotation',
        'subtitle': 'K₁₂ vs K₂₁ with 4 K_rot patterns',
        'p1': 'K₁₂', 'p2': 'K₂₁',
        'grid': None,
        'cases': {
            'A': ('No rotation', 'kpos_x_krot_A.png'),
            'B': ('Symmetric R₁₂=R₂₁=+1', 'kpos_x_krot_B.png'),
            'C': ('Antisymmetric R₁₂=+1,R₂₁=-1', 'kpos_x_krot_C.png'),
            'D': ('One-way R₁₂=+1,R₂₁=0', 'kpos_x_krot_D.png'),
        },
    },
    {
        'folder': 'krot_offdiag',
        'title': 'Sweep 3: Cross-Species Rotation Coupling',
        'subtitle': 'R₁₂ vs R₂₁ (K_pos fixed attractive)',
        'p1': 'R₁₂', 'p2': 'R₂₁',
        'grid': 'krot_offdiag.png',
        'cases': None,
    },
    {
        'folder': 'kpos_diag',
        'title': 'Sweep 4: Self-Cohesion Asymmetry',
        'subtitle': 'K₁₁ vs K₂₂ (K₁₂=K₂₁=0.3, K_rot=0)',
        'p1': 'K₁₁', 'p2': 'K₂₂',
        'grid': 'kpos_diag.png',
        'cases': None,
    },
    {
        'folder': 'krot_diag',
        'title': 'Sweep 5: Self-Rotation',
        'subtitle': 'R₁₁ vs R₂₂ (K_pos fixed attractive)',
        'p1': 'R₁₁', 'p2': 'R₂₂',
        'grid': 'krot_diag.png',
        'cases': None,
    },
]


# ============================================================
# Slide helpers
# ============================================================

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = '2-Species Characterization'
    slide.placeholders[1].text = (
        'Parameter sweep: K_pos and K_rot interaction matrices\n'
        '5-point grid [-1.0, -0.5, 0.0, 0.5, 1.0] per axis'
    )


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=None):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = RGBColor(*color)


def add_grid_image_slide(prs, title, subtitle, grid_path):
    """Slide with the pre-rendered 5x5 screenshot grid image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, 0.3, 0.1, 12, 0.4, title, size=20, bold=True)
    add_text(slide, 0.3, 0.45, 12, 0.3, subtitle, size=13, color=(100, 100, 100))

    if grid_path.exists():
        # Center the image
        img_w, img_h = 8.5, 6.3
        x = (SLIDE_W - img_w) / 2
        slide.shapes.add_picture(str(grid_path), Inches(x), Inches(0.85),
                                 Inches(img_w), Inches(img_h))


def add_video_grid_slide(prs, title, subtitle, p1_label, p2_label,
                         screenshot_dir, video_dir, case=None):
    """Slide with 5x5 grid of embedded videos, axis labels."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Grid layout — fill entire slide with tiny labels
    margin_left = 0.6
    margin_top = 0.25
    cell_w = (SLIDE_W - margin_left) / 5
    cell_h = (SLIDE_H - margin_top) / 5
    vid_size = min(cell_w, cell_h) - 0.02

    # Column labels (top edge) — with parameter name
    for col, p1 in enumerate(GRID_VALUES):
        cx = margin_left + col * cell_w + cell_w / 2 - 0.5
        add_text(slide, cx, 0.0, 1.0, 0.25, f'{p1_label}={p1:+.1f}', size=10, bold=True, color=(200, 0, 0))

    for row, p2 in enumerate(reversed(GRID_VALUES)):
        # Row label (left edge) — with parameter name
        ry = margin_top + row * cell_h + cell_h / 2 - 0.15
        add_text(slide, 0.0, ry, 0.6, 0.3, f'{p2_label}={p2:+.1f}', size=9, bold=True, color=(200, 0, 0))

        for col, p1 in enumerate(GRID_VALUES):
            if case:
                fname = f"{p1:.1f}_{p2:.1f}_{case}"
            else:
                fname = f"{p1:.1f}_{p2:.1f}"

            vid_path = video_dir / f"{fname}.mp4"
            img_path = screenshot_dir / f"{fname}.png"

            x = margin_left + col * cell_w + (cell_w - vid_size) / 2
            y = margin_top + row * cell_h + (cell_h - vid_size) / 2

            if vid_path.exists() and img_path.exists():
                slide.shapes.add_movie(
                    str(vid_path),
                    Inches(x), Inches(y),
                    Inches(vid_size), Inches(vid_size),
                    poster_frame_image=str(img_path)
                )
            elif img_path.exists():
                slide.shapes.add_picture(
                    str(img_path),
                    Inches(x), Inches(y),
                    Inches(vid_size), Inches(vid_size)
                )


# ============================================================
# Main
# ============================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    add_title_slide(prs)

    for sweep in SWEEPS:
        folder = sweep['folder']
        screenshot_dir = RESULTS_DIR / folder / 'screenshots'
        video_dir = RESULTS_DIR / folder / 'videos'

        if sweep['cases']:
            # Sweep 2: one grid per case
            for case, (desc, grid_file) in sweep['cases'].items():
                case_title = f"{sweep['title']} — Case {case}: {desc}"
                grid_path = GRIDS_DIR / grid_file

                # Screenshot grid slide
                add_grid_image_slide(prs, case_title, sweep['subtitle'], grid_path)

                # Video grid slide
                add_video_grid_slide(prs, case_title + ' [Videos]',
                                     sweep['subtitle'],
                                     sweep['p1'], sweep['p2'],
                                     screenshot_dir, video_dir, case=case)
        else:
            grid_path = GRIDS_DIR / sweep['grid']

            # Screenshot grid slide
            add_grid_image_slide(prs, sweep['title'], sweep['subtitle'], grid_path)

            # Video grid slide
            add_video_grid_slide(prs, sweep['title'] + ' [Videos]',
                                 sweep['subtitle'],
                                 sweep['p1'], sweep['p2'],
                                 screenshot_dir, video_dir)

    output_path = OUTPUT_DIR / 'characterization_slides.pptx'
    prs.save(str(output_path))
    print(f"Saved: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == '__main__':
    main()
