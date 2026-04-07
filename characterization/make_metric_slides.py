#!/usr/bin/env python3
"""
Generate organized PowerPoint slides from metric heatmaps.

Each sweep gets:
- A title slide with sweep description
- An overview slide showing all 14 metrics in a grid
- Individual metric slides with explanation

For multi-case sweeps (kpos_x_krot), each metric gets one slide with all 4 cases side-by-side.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


METRICS_DIR = Path(__file__).parent / 'results' / 'metrics'
OUTPUT_DIR = Path(__file__).parent
SLIDE_W = 13.33
SLIDE_H = 7.5


# ============================================================
# Sweep definitions and explanations
# ============================================================

SWEEPS = [
    {
        'folder': 'kpos_x_krot',
        'title': 'Sweep 1: Position × Rotation Cases',
        'subtitle': 'K₁₂ vs K₂₁ with 4 K_rot patterns (A/B/C/D)',
        'description': (
            'Sweeping the off-diagonal entries of K_pos (cross-species '
            'attraction), repeated for 4 K_rot patterns:\n'
            '  A: No rotation\n'
            '  B: Symmetric — collective rotation\n'
            '  C: Antisymmetric — translation\n'
            '  D: One-way — shepherd-like'
        ),
        'example_kpos': [['0.6', 'K₁₂'], ['K₂₁', '0.6']],
        'example_krot_cases': [
            ('A', [['0', '0'], ['0', '0']]),
            ('B', [['0', '+1'], ['+1', '0']]),
            ('C', [['0', '+1'], ['-1', '0']]),
            ('D', [['0', '+1'], ['0', '0']]),
        ],
        'has_overview': False,
    },
    {
        'folder': 'krot_offdiag',
        'title': 'Sweep 2: Cross-Species Rotation Coupling',
        'subtitle': 'R₁₂ vs R₂₁  |  K_pos fixed attractive',
        'description': (
            'Sweeping the off-diagonal entries of K_rot. With a mildly '
            'attractive K_pos baseline, this isolates the effect of '
            'rotation coupling between species pairs.'
        ),
        'example_kpos': [['0.6', '0.3'], ['0.3', '0.6']],
        'example_krot': [['0', 'R₁₂'], ['R₂₁', '0']],
        'has_overview': True,
    },
    {
        'folder': 'kpos_diag',
        'title': 'Sweep 3: Self-Cohesion Asymmetry',
        'subtitle': 'K₁₁ vs K₂₂  |  K₁₂ = K₂₁ = 0.3, K_rot = 0',
        'description': (
            'Sweeping the diagonal entries of K_pos. K₁₁ is species 1 '
            'self-cohesion, K₂₂ is species 2 self-cohesion. Cross-attraction '
            'fixed at 0.3. Reveals what happens when species have different '
            'internal cohesion strengths.'
        ),
        'example_kpos': [['K₁₁', '0.3'], ['0.3', 'K₂₂']],
        'example_krot': [['0', '0'], ['0', '0']],
        'has_overview': True,
    },
    {
        'folder': 'krot_diag',
        'title': 'Sweep 4: Self-Rotation',
        'subtitle': 'R₁₁ vs R₂₂  |  K_pos fixed attractive, R₁₂ = R₂₁ = 0',
        'description': (
            'Sweeping the diagonal entries of K_rot. Each diagonal entry '
            'controls how a species rotates around its own centroid '
            '(self-rotation), independent of the other species.'
        ),
        'example_kpos': [['0.6', '0.3'], ['0.3', '0.6']],
        'example_krot': [['R₁₁', '0'], ['0', 'R₂₂']],
        'has_overview': True,
    },
]


# ============================================================
# Metric definitions (in display order)
# ============================================================

METRIC_GROUPS = [
    {
        'name': 'Spatial Spread',
        'metrics': [
            ('max_d1', 'Max Distance (Group 1)',
             'Maximum pairwise distance within species 1. Large values mean group 1 is spread across the arena; small values mean it is clustered.'),
            ('max_d2', 'Max Distance (Group 2)',
             'Maximum pairwise distance within species 2. Same interpretation as group 1.'),
            ('centroid_dist', 'Centroid Distance',
             'Distance between the centers of mass of group 1 and group 2. High values = groups separated; low values = concentric or overlapping.'),
        ],
    },
    {
        'name': 'Speed & Energy',
        'metrics': [
            ('avg_speed', 'Avg Speed (All)',
             'Average particle speed across both groups. High speeds indicate active dynamics (chase, repulsion); low speeds indicate stable equilibrium.'),
            ('avg_speed1', 'Avg Speed (Group 1)',
             'Average speed of species 1 only. Asymmetry between this and group 2 reveals which species is more active.'),
            ('avg_speed2', 'Avg Speed (Group 2)',
             'Average speed of species 2 only.'),
            ('KE', 'Kinetic Energy',
             'Mean squared velocity (½v²) — proxy for activity level. High KE marks chaotic or fast-moving regions of parameter space.'),
        ],
    },
    {
        'name': 'Mean Square Displacement',
        'metrics': [
            ('MSD', 'MSD (All)',
             'Mean squared displacement of all particles from their initial positions during the measurement window. Measures how far particles travel — distinguishes diffusive (high MSD) from confined (low MSD) regimes.'),
            ('MSD1', 'MSD (Group 1)',
             'MSD for species 1 only.'),
            ('MSD2', 'MSD (Group 2)',
             'MSD for species 2 only.'),
        ],
    },
    {
        'name': 'Inter-Particle Spacing',
        'metrics': [
            ('spacing_all', 'Avg Spacing (All)',
             'Mean pairwise distance between all particles regardless of species. Indicates global density.'),
            ('spacing_same', 'Avg Spacing (Same Group)',
             'Average of the within-group spacings (mean of group 1 and group 2). Shows how compact each species is internally.'),
            ('spacing1', 'Avg Spacing (Group 1)',
             'Mean pairwise distance within species 1.'),
            ('spacing2', 'Avg Spacing (Group 2)',
             'Mean pairwise distance within species 2.'),
        ],
    },
]


# ============================================================
# Slide builders
# ============================================================

def add_text(slide, x, y, w, h, text, size=14, bold=False, color=None, align=None):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = RGBColor(*color)
    if align:
        p.alignment = align
    return txBox


def draw_matrix(slide, x, y, label, values, cell_size=0.55, label_size=14):
    """Draw a 2x2 matrix as text labels with cell borders."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Pt as _Pt

    # Label above matrix
    add_text(slide, x, y, cell_size * 2 + 0.5, 0.3, label,
             size=label_size, bold=True, color=(30, 30, 80))
    y += 0.32

    # Brackets and cells
    n = 2
    for i in range(n):
        for j in range(n):
            cx = x + j * cell_size
            cy = y + i * cell_size
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(cx), Inches(cy),
                Inches(cell_size - 0.05), Inches(cell_size - 0.05))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(245, 245, 250)
            shape.line.color.rgb = RGBColor(80, 80, 120)
            shape.line.width = Pt(1)

            tf = shape.text_frame
            tf.margin_top = Pt(2)
            tf.margin_bottom = Pt(2)
            tf.margin_left = Pt(2)
            tf.margin_right = Pt(2)
            p = tf.paragraphs[0]
            p.alignment = 2  # Center
            p.text = str(values[i][j])
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = RGBColor(30, 30, 80)


def add_section_title_slide(prs, sweep):
    """Title page for each sweep with description and example matrices."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    add_text(slide, 0.5, 0.6, 12.3, 0.6, sweep['title'],
             size=30, bold=True, color=(30, 30, 80))
    add_text(slide, 0.5, 1.3, 12.3, 0.4, sweep['subtitle'],
             size=16, color=(100, 100, 100))

    # Description (left side)
    add_text(slide, 0.5, 2.2, 7.0, 0.4,
             'Description', size=16, bold=True, color=(30, 30, 80))
    add_text(slide, 0.5, 2.7, 7.0, 4.5, sweep['description'],
             size=14, color=(60, 60, 60))

    # Example matrices (right side)
    add_text(slide, 8.0, 2.2, 5.0, 0.4,
             'Example Matrices', size=16, bold=True, color=(30, 30, 80))

    if 'example_krot_cases' in sweep:
        # Multi-case sweep — show K_pos once, then 4 K_rot cases
        draw_matrix(slide, 8.2, 2.7, 'K_pos:', sweep['example_kpos'])

        # 4 K_rot cases in a 2x2 grid
        case_x_start = 10.5
        case_y_start = 2.7
        for idx, (case_name, mat) in enumerate(sweep['example_krot_cases']):
            row = idx // 2
            col = idx % 2
            cx = case_x_start + col * 1.4
            cy = case_y_start + row * 1.85
            draw_matrix(slide, cx, cy, f'K_rot ({case_name}):', mat,
                        cell_size=0.45, label_size=11)
    else:
        # Single matrix case
        draw_matrix(slide, 8.2, 2.7, 'K_pos:', sweep['example_kpos'])
        draw_matrix(slide, 11.0, 2.7, 'K_rot:', sweep['example_krot'])


def add_overview_slide(prs, sweep, overview_path):
    """Slide showing all metrics in one big image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text(slide, 0.3, 0.1, 12, 0.4, f"{sweep['title']} — All Metrics Overview",
             size=18, bold=True)
    add_text(slide, 0.3, 0.5, 12, 0.3, sweep['subtitle'],
             size=12, color=(100, 100, 100))

    if overview_path.exists():
        # Maximize image
        img_w, img_h = 11.5, 6.4
        x = (SLIDE_W - img_w) / 2
        slide.shapes.add_picture(str(overview_path), Inches(x), Inches(0.95),
                                 Inches(img_w), Inches(img_h))


def add_metric_slide(prs, sweep, metric_key, metric_label, explanation, image_path):
    """Slide for one metric: heatmap on left, explanation on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    add_text(slide, 0.3, 0.1, 12.5, 0.4,
             f"{sweep['title']} — {metric_label}",
             size=18, bold=True)
    add_text(slide, 0.3, 0.5, 12.5, 0.3, sweep['subtitle'],
             size=11, color=(100, 100, 100))

    # Heatmap (left side)
    if image_path.exists():
        img_w, img_h = 6.5, 6.0
        slide.shapes.add_picture(str(image_path), Inches(0.3), Inches(1.0),
                                 Inches(img_w), Inches(img_h))

    # Explanation (right side)
    add_text(slide, 7.2, 1.2, 6.0, 0.4,
             'Description', size=14, bold=True, color=(30, 30, 80))
    add_text(slide, 7.2, 1.7, 6.0, 4.5, explanation, size=13, color=(60, 60, 60))

    # Reading guide (right side, lower)
    add_text(slide, 7.2, 5.5, 6.0, 0.3,
             'How to read the heatmap', size=12, bold=True, color=(30, 30, 80))
    p1 = sweep['subtitle'].split('|')[0].strip().split('vs')
    p1_name = p1[0].strip() if len(p1) > 0 else 'param1'
    p2_name = p1[1].strip() if len(p1) > 1 else 'param2'
    reading = (f'X-axis: {p1_name}    Y-axis: {p2_name}\n'
               'Color: metric value (see colorbar)\n'
               'Each cell = averaged over the last 20% of one simulation')
    add_text(slide, 7.2, 5.8, 6.0, 1.2, reading, size=11, color=(100, 100, 100))


def add_multicase_metric_slide(prs, sweep, metric_key, metric_label, explanation,
                                metrics_dir):
    """For multi-case sweeps: 4 heatmaps in 2x2 grid + explanation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text(slide, 0.3, 0.1, 12.5, 0.4,
             f"{sweep['title']} — {metric_label}",
             size=18, bold=True)
    add_text(slide, 0.3, 0.5, 12.5, 0.3, sweep['subtitle'],
             size=11, color=(100, 100, 100))

    # The metric file shows all 4 cases side by side already
    image_path = metrics_dir / sweep['folder'] / f"{metric_key}.png"
    if image_path.exists():
        img_w, img_h = 12.5, 4.0
        slide.shapes.add_picture(str(image_path), Inches(0.4), Inches(1.0),
                                 Inches(img_w), Inches(img_h))

    # Explanation at bottom
    add_text(slide, 0.5, 5.4, 12.3, 0.4,
             'Description', size=14, bold=True, color=(30, 30, 80))
    add_text(slide, 0.5, 5.85, 12.3, 1.5, explanation, size=12, color=(60, 60, 60))


def add_metric_section_divider(prs, group_name):
    """Section divider between metric groups."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, 0.5, 3.0, 12.3, 1.0, group_name,
             size=44, bold=True, color=(30, 30, 80))


# ============================================================
# Main
# ============================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # === Cover slide ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, 0.5, 2.5, 12.3, 1.0,
             '2-Species Quantitative Characterization',
             size=40, bold=True, color=(30, 30, 80))
    add_text(slide, 0.5, 3.7, 12.3, 0.6,
             'Metric heatmaps across 5 parameter sweeps',
             size=22, color=(100, 100, 100))
    add_text(slide, 0.5, 4.5, 12.3, 0.4,
             'Grid: 11×11, range [-1.0, 1.0], measured over last 20% of simulation',
             size=14, color=(120, 120, 120))

    # === Each sweep ===
    for sweep in SWEEPS:
        sweep_dir = METRICS_DIR / sweep['folder']
        if not sweep_dir.exists():
            print(f"Skipping {sweep['folder']} — no results")
            continue

        # 1. Section title
        add_section_title_slide(prs, sweep)

        # 2. Overview (if available)
        if sweep['has_overview']:
            overview_path = sweep_dir / 'overview.png'
            add_overview_slide(prs, sweep, overview_path)

        # 3. Metric slides grouped by category
        for group in METRIC_GROUPS:
            for metric_key, metric_label, explanation in group['metrics']:
                if sweep['has_overview']:
                    image_path = sweep_dir / f"{metric_key}.png"
                    if image_path.exists():
                        add_metric_slide(prs, sweep, metric_key, metric_label,
                                         explanation, image_path)
                else:
                    # Multi-case sweep
                    image_path = sweep_dir / f"{metric_key}.png"
                    if image_path.exists():
                        add_multicase_metric_slide(prs, sweep, metric_key,
                                                   metric_label, explanation,
                                                   METRICS_DIR)

    output_path = OUTPUT_DIR / 'metric_slides.pptx'
    prs.save(str(output_path))
    print(f"Saved: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == '__main__':
    main()
